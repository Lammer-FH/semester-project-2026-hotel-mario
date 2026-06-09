from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ────────────────────────────────────────────────────────────────
PRIMARY   = RGBColor(0x3E, 0x80, 0xC8)   # Ionic blue
DARK      = RGBColor(0x1A, 0x1A, 0x2E)   # near-black
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF4, 0xF7, 0xFD)
ACCENT    = RGBColor(0xE8, 0x5D, 0x04)   # orange accent
SUBTLE    = RGBColor(0x6B, 0x7A, 0x99)
CODE_BG   = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG   = RGBColor(0xCD, 0xD6, 0xF4)

W = Inches(13.33)   # 16:9 widescreen
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def add_bullet_box(slide, items, x, y, w, h,
                   size=16, color=DARK, bullet="▸ ", bold_first=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        if bold_first:
            run.font.bold = True
            bold_first = False
    return tb

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.35), DARK)
    add_text(slide, title,
             Inches(0.5), Inches(0.18), Inches(10), Inches(0.7),
             size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.85), Inches(10), Inches(0.45),
                 size=16, color=PRIMARY, italic=False)
    # accent stripe
    add_rect(slide, 0, Inches(1.35), W, Inches(0.06), PRIMARY)

def slide_bg(slide, color=LIGHT_BG):
    add_rect(slide, 0, 0, W, H, color)

def tag_pill(slide, text, x, y, w=Inches(1.6), h=Inches(0.38), bg=PRIMARY):
    add_rect(slide, x, y, w, h, bg)
    add_text(slide, text, x, y, w, h,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def code_block(slide, lines, x, y, w, h):
    add_rect(slide, x, y, w, h, CODE_BG)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(11)
        run.font.color.rgb = CODE_FG
        run.font.name = "Courier New"


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK)
add_rect(s, 0, 0, Inches(0.18), H, PRIMARY)          # left accent bar

add_text(s, "Boutique Hotel Technikum",
         Inches(0.55), Inches(1.6), Inches(12), Inches(1.1),
         size=44, bold=True, color=WHITE)
add_text(s, "A full-stack hotel booking app",
         Inches(0.55), Inches(2.65), Inches(9), Inches(0.7),
         size=22, color=PRIMARY)
add_text(s, "Group AE  ·  Semester Project 2026  ·  MSE25",
         Inches(0.55), Inches(3.35), Inches(9), Inches(0.55),
         size=16, color=SUBTLE, italic=True)

# tech chips row
chips = ["Ionic Vue 7", "Spring Boot 4", "MySQL 8", "Clean Architecture", "Pinia", "Testcontainers"]
cx = Inches(0.55)
for c in chips:
    tag_pill(s, c, cx, Inches(4.4), Inches(1.75), Inches(0.38))
    cx += Inches(1.9)

add_text(s, "15-min presentation  ·  June 2026",
         Inches(0.55), Inches(6.7), Inches(8), Inches(0.5),
         size=13, color=SUBTLE)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Live Demo Guide (speaker notes aid, minimal text)
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Live Demo", "Walk the user journey end-to-end  (~5 min)")

steps = [
    ("1  Landing Page",   "Home screen — hotel overview, slideshow, CTA button"),
    ("2  Room Selection", "Filter by date range (two-step date picker), price, available-only toggle"),
    ("3  Pagination",     "Step through pages — server-side, 5 rooms per page"),
    ("4  Book a Room",    "Select available room → form with validation → review screen"),
    ("5  Confirmation",   "Show confirmation card: booking ID, room image, hotel details, print to A4"),
    ("6  Error path",     "Attempt to book an already-taken room → 409 'Room no longer available'"),
]

y = Inches(1.6)
for title, desc in steps:
    add_rect(s, Inches(0.45), y, Inches(0.07), Inches(0.42), PRIMARY)
    add_text(s, title, Inches(0.65), y - Inches(0.03), Inches(3.2), Inches(0.5),
             size=15, bold=True, color=DARK)
    add_text(s, desc,  Inches(3.85), y - Inches(0.03), Inches(9.0), Inches(0.5),
             size=14, color=SUBTLE)
    y += Inches(0.78)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — AI Workflow (simple, one slide)
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "AI Workflow", "Clean Architecture Scaffolding from API Spec  ·  ~5 min")

# Left — the prompt
add_text(s, "What we gave AI",
         Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.45),
         size=18, bold=True, color=DARK)
add_bullet_box(s, [
    "Our ER diagram (ROOM, EXTRA, BOOKING, ROOM_EXTRA)",
    "The OpenAPI spec we wrote (endpoints, DTOs, pagination shape)",
    "Architecture decision: component packages, v1 versioning, interface/impl split, Lombok",
], Inches(0.5), Inches(2.1), Inches(5.8), Inches(1.3), size=13, color=DARK)

add_text(s, "What AI generated (one pass)",
         Inches(0.5), Inches(3.5), Inches(5.8), Inches(0.45),
         size=18, bold=True, color=PRIMARY)
add_bullet_box(s, [
    "Entities, JPA repositories, service interfaces + impls, mappers",
    "@EmailsMatch and @CheckOutAfterCheckIn as proper class-level constraint annotations",
    "GlobalExceptionHandler returning structured {status, message, errors[]}",
    "13 E2E tests (Testcontainers + MockMvc) covering overlap, 409, pagination",
], Inches(0.5), Inches(4.0), Inches(5.8), Inches(2.1), size=13, color=DARK)

add_text(s, "What we changed",
         Inches(0.5), Inches(6.2), Inches(5.8), Inches(0.45),
         size=14, bold=True, color=ACCENT)
add_text(s, "Minor corrections only — one wrong Spring Boot 4 import, one overly broad exception catch.",
         Inches(0.5), Inches(6.55), Inches(5.8), Inches(0.45),
         size=12, color=DARK, italic=True)

# Right — example code snippet
add_text(s, "Example: cross-field constraint annotation",
         Inches(6.6), Inches(1.6), Inches(6.3), Inches(0.45),
         size=16, bold=True, color=DARK)
code_block(s, [
    "@EmailsMatch                           // class-level",
    "@CheckOutAfterCheckIn                  // class-level",
    "@Data @Builder @NoArgsConstructor",
    "public class BookingRequestDto {",
    "  @NotBlank @Email  String email;",
    "  @NotBlank         String emailConfirmation;",
    "  @FutureOrPresent  LocalDate checkIn;",
    "                    LocalDate checkOut;",
    "  // ...",
    "}",
    "",
    "// Validator generated correctly first attempt:",
    "public class EmailsMatchValidator implements",
    "    ConstraintValidator<EmailsMatch, BookingRequestDto> {",
    "  public boolean isValid(BookingRequestDto dto, ...) {",
    "    return dto.getEmail().equals(dto.getEmailConfirmation());",
    "  }",
    "}",
], Inches(6.6), Inches(2.1), Inches(6.4), Inches(4.5))

add_text(s, "Human decision: we defined the architecture and validation rules — AI filled in the boilerplate correctly.",
         Inches(6.6), Inches(6.65), Inches(6.3), Inches(0.55),
         size=11, bold=True, color=ACCENT)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — What surprised us
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "What Surprised Us", "One thing easier than expected · one thing it got wrong")

# Left column — positive
add_rect(s, Inches(0.4), Inches(1.55), Inches(5.9), Inches(5.5), WHITE)
add_rect(s, Inches(0.4), Inches(1.55), Inches(5.9), Inches(0.55), PRIMARY)
add_text(s, "✓  Easier than expected",
         Inches(0.55), Inches(1.6), Inches(5.6), Inches(0.45),
         size=16, bold=True, color=WHITE)

positives = [
    ("Full-stack scaffolding from architecture doc",
     "We handed AI our ER diagram and API spec. It generated the entire clean-architecture backend skeleton — entities, repos, services, mappers, DTOs, controllers — in one pass. Only minor corrections needed."),
    ("Cross-field validation annotations",
     "@EmailsMatch and @CheckOutAfterCheckIn as proper class-level constraint annotations with validator classes — correct on the first attempt. Would have taken us 30 min to research and write manually."),
    ("Test coverage from specs",
     "Given endpoint descriptions, AI wrote 13 meaningful E2E tests (Testcontainers + MockMvc) including overlap checks and 409 conflict cases we hadn't thought of."),
]
y = Inches(2.2)
for title, detail in positives:
    add_text(s, title, Inches(0.6), y, Inches(5.5), Inches(0.4),
             size=13, bold=True, color=PRIMARY)
    add_text(s, detail, Inches(0.6), y + Inches(0.38), Inches(5.5), Inches(0.85),
             size=12, color=DARK)
    y += Inches(1.4)

# Right column — negative
add_rect(s, Inches(7.0), Inches(1.55), Inches(5.9), Inches(5.5), WHITE)
add_rect(s, Inches(7.0), Inches(1.55), Inches(5.9), Inches(0.55), ACCENT)
add_text(s, "✗  Where it led us astray",
         Inches(7.15), Inches(1.6), Inches(5.6), Inches(0.45),
         size=16, bold=True, color=WHITE)

negatives = [
    ("Fake pagination (peer review catch)",
     "AI generated fetchRooms() hardcoded to getRooms(0, 20) — all rooms, no real pages. Pagination buttons existed but were cosmetic. Only caught during peer review. Root cause: we didn't specify 'real server-side pagination' explicitly."),
    ("Wrong TypeScript fix loops",
     "The DatetimeHighlightStyle type error: AI tried adding border: 'none' (wrong property) instead of removing the return type annotation. Took two prompts to converge because we didn't explain the union type constraint upfront."),
    ("imageUrl mismatch",
     "Backend DataSeeder used filename-based paths (/images/rooms/standard.jpg) but the actual public files are numbered (/images/rooms/1.jpg). AI never caught this — we found it visually. Lesson: AI can't verify runtime asset resolution."),
]
y = Inches(2.2)
for title, detail in negatives:
    add_text(s, title, Inches(7.2), y, Inches(5.5), Inches(0.4),
             size=13, bold=True, color=ACCENT)
    add_text(s, detail, Inches(7.2), y + Inches(0.38), Inches(5.5), Inches(0.85),
             size=12, color=DARK)
    y += Inches(1.4)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Architecture overview (clean architecture diagram as text)
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Architecture at a Glance", "Clean Architecture · Atomic Design · Pinia stores")

# Backend box
add_rect(s, Inches(0.4), Inches(1.55), Inches(5.9), Inches(5.5), WHITE)
add_rect(s, Inches(0.4), Inches(1.55), Inches(5.9), Inches(0.45), DARK)
add_text(s, "Backend — Spring Boot 4",
         Inches(0.55), Inches(1.58), Inches(5.6), Inches(0.4),
         size=15, bold=True, color=WHITE)

be_layers = [
    ("config/",              "HotelConfig, DataSeeder, WebConfig"),
    ("common/",              "GlobalExceptionHandler, PageResult, exceptions"),
    ("components/rooms/",    "api/v1 → service → repository → model"),
    ("components/bookings/", "api/v1 → service → repository → model"),
    ("  api/dtos/v1/",       "RoomDto, BookingRequestDto, BookingConfirmationDto …"),
    ("  api/mapper/v1/",     "RoomMapper, BookingMapper (interface + impl)"),
    ("  service/",           "RoomService, RoomAvailabilityService (interface + impl)"),
    ("  repository/",        "JPA: RoomRepository, RoomAvailabilityRepository"),
]
y = Inches(2.1)
for pkg, desc in be_layers:
    add_text(s, pkg, Inches(0.6),  y, Inches(2.1), Inches(0.38),
             size=11, bold=True, color=PRIMARY)
    add_text(s, desc, Inches(2.7), y, Inches(3.4), Inches(0.38),
             size=11, color=DARK)
    y += Inches(0.42)

# Frontend box
add_rect(s, Inches(7.0), Inches(1.55), Inches(5.9), Inches(5.5), WHITE)
add_rect(s, Inches(7.0), Inches(1.55), Inches(5.9), Inches(0.45), DARK)
add_text(s, "Frontend — Ionic Vue 7",
         Inches(7.15), Inches(1.58), Inches(5.6), Inches(0.4),
         size=15, bold=True, color=WHITE)

fe_layers = [
    ("views/",                  "RoomSelectionView, BookingView, BookingReviewView, BookingDetailsView"),
    ("components/atoms/",       "ExtraChip, AvailabilityBadge, DetailRow"),
    ("components/molecules/",   "RoomCard, BookingForm, DatePickerModal, ImageSlider, SharedHeader/Footer"),
    ("components/organisms/",   "RoomList, FilterBar"),
    ("stores/",                 "useRoomStore (rooms, pagination), useFilterStore (dates, price), useBookingStore"),
    ("services/api.ts",         "Typed DTOs, getRooms(), checkAvailability(), bookRoom()"),
    ("router/",                 "Vue Router — /home, /room, /booking/:id, /booking/details"),
]
y = Inches(2.1)
for pkg, desc in fe_layers:
    add_text(s, pkg, Inches(7.2),  y, Inches(2.2), Inches(0.38),
             size=11, bold=True, color=PRIMARY)
    add_text(s, desc, Inches(9.4), y, Inches(3.3), Inches(0.38),
             size=11, color=DARK)
    y += Inches(0.52)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Q&A
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK)
add_rect(s, 0, 0, Inches(0.18), H, PRIMARY)

add_text(s, "Questions?",
         Inches(0.55), Inches(1.8), Inches(12), Inches(1.2),
         size=52, bold=True, color=WHITE)
add_text(s, "One question per team — AI usage, architecture, things that went wrong",
         Inches(0.55), Inches(3.0), Inches(10), Inches(0.7),
         size=20, color=PRIMARY)

topics = [
    "How did you decide what to verify vs. just accept from AI?",
    "What would you prompt differently next time?",
    "Where did clean architecture boundaries actually help you?",
    "Did AI-generated tests catch real bugs?",
]
y = Inches(3.85)
for t in topics:
    add_rect(s, Inches(0.55), y, Inches(0.06), Inches(0.32), ACCENT)
    add_text(s, t, Inches(0.75), y - Inches(0.04), Inches(11.5), Inches(0.42),
             size=14, color=SUBTLE)
    y += Inches(0.52)

add_text(s, "Group AE  ·  Michael Auß  ·  Boutique Hotel Technikum",
         Inches(0.55), Inches(6.7), Inches(11), Inches(0.5),
         size=13, color=SUBTLE, italic=True)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8b — Peer Review 2 Response
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s, "Peer Review 2 — Feedback & Response", "Component complexity · code duplication · atomic design")

add_text(s, "What the reviewer flagged",
         Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.45),
         size=18, bold=True, color=DARK)
add_bullet_box(s, [
    "Duplicated template logic across views (same label/value rows repeated by hand)",
    "Duplicated slideshow code in home.vue and about.vue (~80 lines each)",
    "No consistent reuse of existing atoms (ExtraChip used in RoomCard but not in confirmation view)",
], Inches(0.5), Inches(2.1), Inches(5.8), Inches(1.6), size=13, color=DARK)

add_text(s, "What we changed (with AI)",
         Inches(0.5), Inches(3.8), Inches(5.8), Inches(0.45),
         size=18, bold=True, color=PRIMARY)
changes = [
    ("atoms/DetailRow.vue",   "New atom: <strong>Label:</strong> Value row — replaced 16 ion-item blocks\nacross BookingDetailsView + BookingReviewView"),
    ("molecules/ImageSlider.vue", "New molecule: reusable slideshow — removed ~160 lines of duplicated\nscript + template from home.vue and about.vue"),
    ("ExtraChip (extended)",  "Applied existing atom in BookingDetailsView to match\nRoomCard — consistent chip rendering throughout"),
]
y = Inches(4.3)
for comp, desc in changes:
    add_rect(s, Inches(0.5), y + Inches(0.06), Inches(0.06), Inches(0.3), PRIMARY)
    add_text(s, comp, Inches(0.7), y, Inches(2.8), Inches(0.42),
             size=12, bold=True, color=PRIMARY)
    add_text(s, desc, Inches(3.5), y, Inches(3.0), Inches(0.55),
             size=11, color=DARK)
    y += Inches(0.75)

# Right panel — before/after component diagram
add_rect(s, Inches(7.0), Inches(1.55), Inches(5.9), Inches(5.5), WHITE)
add_rect(s, Inches(7.0), Inches(1.55), Inches(5.9), Inches(0.45), DARK)
add_text(s, "Atomic Design structure (after)",
         Inches(7.15), Inches(1.58), Inches(5.6), Inches(0.4),
         size=14, bold=True, color=WHITE)

atoms_list = [
    ("atoms/",      "ExtraChip  ·  AvailabilityBadge  ·  DetailRow"),
    ("molecules/",  "RoomCard  ·  BookingForm  ·  DatePickerModal"),
    ("",            "ImageSlider  ·  SharedHeader  ·  SharedFooter"),
    ("organisms/",  "RoomList  ·  FilterBar"),
]
y = Inches(2.1)
for label, items in atoms_list:
    if label:
        add_text(s, label, Inches(7.2), y, Inches(1.5), Inches(0.38),
                 size=12, bold=True, color=PRIMARY)
    add_text(s, items, Inches(8.7), y, Inches(3.9), Inches(0.38),
             size=12, color=DARK)
    y += Inches(0.44)

add_text(s, "What we scoped out (not actioned)",
         Inches(7.2), Inches(4.1), Inches(5.5), Inches(0.42),
         size=14, bold=True, color=ACCENT)
add_bullet_box(s, [
    "HotelConfig → HotelDto cross-layer dep. — non-functional, deferred",
    "FilterBar untyped props — low priority, no TypeScript errors",
], Inches(7.2), Inches(4.55), Inches(5.5), Inches(1.0), size=12, color=DARK)

add_text(s, "Human decision: scoped to high-impact duplications only — explicit choice not to over-refactor",
         Inches(7.2), Inches(5.65), Inches(5.5), Inches(0.7),
         size=11, bold=True, color=SUBTLE, italic=True)


# ── Save ───────────────────────────────────────────────────────────────────
out = "/Users/michael/Projekte/semester-project-2026-hotel-mario/presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
